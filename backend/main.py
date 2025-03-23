from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import fitz  # PyMuPDF for PDF
import docx  # python-docx for DOC & DOCX
from pptx import Presentation  # pptx for PowerPoint
import nltk
from nltk.tokenize import sent_tokenize
from nltk.corpus import stopwords
from collections import Counter
import string

# Download necessary NLTK data
nltk.download("punkt")
nltk.download("stopwords")

app = FastAPI()

# Allow CORS for frontend access
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins (Change in production)
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def extract_text_from_pdf(file):
    """Extracts text from a PDF file using PyMuPDF."""
    try:
        doc = fitz.open(stream=file, filetype="pdf")
        text = "\n".join([page.get_text("text") for page in doc])
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text from PDF: {str(e)}")

def extract_text_from_docx(file):
    """Extracts text from a DOCX file."""
    try:
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text from DOCX: {str(e)}")

def extract_text_from_pptx(file):
    """Extracts text from a PowerPoint (PPTX) file."""
    try:
        prs = Presentation(file)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text from PPTX: {str(e)}")

def generate_summary(text):
    """Generates a summary from the given text using NLTK."""
    sentences = sent_tokenize(text)
    words = text.lower().translate(str.maketrans("", "", string.punctuation)).split()
    
    stop_words = set(stopwords.words("english"))
    filtered_words = [word for word in words if word not in stop_words]
    
    word_freq = Counter(filtered_words)
    sentence_scores = {sentence: sum(word_freq.get(word.lower(), 0) for word in sentence.split()) for sentence in sentences}
    
    top_sentences = sorted(sentence_scores, key=sentence_scores.get, reverse=True)[:3]  # Top 3 sentences
    return " ".join(top_sentences)

@app.post("/summarize/")
async def summarize_file(file: UploadFile = File(...)):
    """Handles file upload and generates a summary."""
    try:
        content = await file.read()

        # Identify file type and extract text
        if file.filename.endswith(".pdf"):
            text = extract_text_from_pdf(content)
        elif file.filename.endswith(".docx"):
            text = extract_text_from_docx(file.file)
        elif file.filename.endswith(".pptx"):
            text = extract_text_from_pptx(file.file)
        elif file.filename.endswith(".txt"):
            text = content.decode("utf-8")
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format. Supported: PDF, DOCX, PPTX, TXT")

        if not text.strip():
            return {"summary": "No content found in the document."}

        summary = generate_summary(text)
        return {"summary": summary, "filename": file.filename}

    except Exception as e:
        return {"error": str(e)}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)

